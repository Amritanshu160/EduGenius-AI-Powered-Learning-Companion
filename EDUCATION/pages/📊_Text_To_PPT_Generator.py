from google import genai
import os
import json
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import io
import streamlit as st
from dotenv import load_dotenv

# Load API Key from environment variables
load_dotenv()


def generate_presentation_content(topic, num_slides, style):
    client = genai.Client(api_key=os.getenv("GOOGLE_API_KEY"))

    try:
        prompt = f"""
        Create a detailed PowerPoint presentation outline on the topic: "{topic}".
        
        The presentation should have {num_slides} slides, including a title slide and a conclusion/summary slide.
        The style should be {style}.
        
        Format your response as a JSON array of slide objects. Each slide object should have the following structure:
        {{
            "title": "Slide Title",
            "content": ["Bullet point 1", "Bullet point 2", "Bullet point 3"],
            "notes": "Optional presenter notes or additional context"
        }}
        
        The first slide should be a title slide with a compelling title, subtitle, and introduction.
        The last slide should summarize key points.
        
        For the remaining slides, create a logical flow with clear headings and concise bullet points.
        
        Important: Return ONLY the JSON array with no additional explanation or text.
        """
        response = client.models.generate_content(
            model = "gemini-2.5-flash",
            contents = prompt
        )
        content_text = response.text

        if "```json" in content_text:
            content_text = content_text.split("```json")[1].split("```")[0].strip()
        elif "```" in content_text:
            content_text = content_text.split("```")[1].strip()

        presentation_content = json.loads(content_text)
        return presentation_content

    except Exception as e:
        print(f"Error generating presentation content: {str(e)}")
        raise

def create_presentation(content):
    prs = Presentation()
    title_slide_layout = prs.slide_layouts[0]
    content_slide_layout = prs.slide_layouts[1]

    # Font size configurations
    TITLE_FONT_SIZE = Pt(36)
    INTRO_SUBTITLE_FONT_SIZE = Pt(18)  # Smaller for intro slide
    CONTENT_FONT_SIZE = Pt(16)
    BULLET_SPACING = Pt(8)

    for i, slide_data in enumerate(content):
        if i == 0:  # Title/intro slide
            slide = prs.slides.add_slide(title_slide_layout)
            title = slide.shapes.title
            subtitle = slide.placeholders[1]
            
            # Set title
            title.text = slide_data["title"]
            title.text_frame.paragraphs[0].font.size = TITLE_FONT_SIZE
            
            # Handle subtitle content with overflow protection
            if slide_data.get("content"):
                subtitle.text = ""  # Clear default text
                text_frame = subtitle.text_frame
                text_frame.word_wrap = True
                text_frame.auto_size = None  # Disable auto-size
                
                # Add each line separately with controlled formatting
                for line in slide_data["content"]:
                    p = text_frame.add_paragraph()
                    p.text = line
                    p.font.size = INTRO_SUBTITLE_FONT_SIZE
                    p.space_after = BULLET_SPACING
                    
                    # Shorten very long lines (optional)
                    if len(line) > 100:
                        p.text = line[:97] + "..."  # Truncate with ellipsis

        else:  # Content slides
            slide = prs.slides.add_slide(content_slide_layout)
            title = slide.shapes.title
            content_shape = slide.placeholders[1]
            
            title.text = slide_data["title"]
            title.text_frame.paragraphs[0].font.size = TITLE_FONT_SIZE
            
            text_frame = content_shape.text_frame
            text_frame.clear()
            text_frame.word_wrap = True
            
            for point in slide_data.get("content", []):
                p = text_frame.add_paragraph()
                p.text = point
                p.level = 0
                p.space_after = BULLET_SPACING
                p.font.size = CONTENT_FONT_SIZE
                
                if len(slide_data.get("content", [])) > 6:
                    p.font.size = Pt(14)

        if slide_data.get("notes"):
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = slide_data["notes"]

    ppt_bytes = io.BytesIO()
    prs.save(ppt_bytes)
    ppt_bytes.seek(0)
    return ppt_bytes.getvalue()

# ✅ Streamlit UI starts here
st.set_page_config(page_title="PPT Generator", layout="wide",page_icon="📊")
st.title("📊 AI Presentation Generator")

with st.form("presentation_form"):
    topic = st.text_input("Enter Presentation Topic")
    num_slides = st.number_input("Number of Slides", min_value=3, max_value=50, value=5)
    style = st.selectbox("Choose Presentation Style", ["Professional", "Academic", "Creative", "Minimal"])
    submitted = st.form_submit_button("Generate Presentation")

if submitted:
    try:
        st.info("Generating content using Gemini...")
        content = generate_presentation_content(topic, num_slides, style)
        pptx_bytes = create_presentation(content)
        st.success("Presentation generated successfully!")
        st.download_button(
            label="📥 Download Presentation",
            data=pptx_bytes,
            file_name=f"{topic.replace(' ', '_')}.pptx",
            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
        )
    except Exception as e:
        st.error(f"An error occurred: {str(e)}")
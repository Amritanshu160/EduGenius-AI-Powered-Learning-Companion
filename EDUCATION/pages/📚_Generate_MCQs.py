import streamlit as st
from PyPDF2 import PdfReader
from docx import Document
from pptx import Presentation
from google import genai
import os
from youtube_transcript_api import YouTubeTranscriptApi
import tempfile
from dotenv import load_dotenv

load_dotenv()

client = genai.Client(api_key=os.getenv("GOOGLE_API_KEY"))


# Extract text from uploaded files
def extract_text_from_files(files):
    text = ""
    for file in files:
        if file.name.endswith(".pdf"):
            pdf_reader = PdfReader(file)
            for page in pdf_reader.pages:
                text += page.extract_text()
        elif file.name.endswith(".docx"):
            doc = Document(file)
            for para in doc.paragraphs:
                text += para.text + "\n"
        elif file.name.endswith((".pptx", ".ppt")):
            presentation = Presentation(file)
            for slide in presentation.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
        else:
            st.warning(f"Unsupported file format: {file.name}")
    return text

# Generate MCQs using Gemini
def generate_mcqs(text):
    prompt = f"""
    You are an expert educator and exam question setter.

    From the following educational content, generate a comprehensive set of **multiple choice questions (MCQs)** that fully cover **all the important facts, concepts, definitions, and key information**.

    - Do **not limit** the number of questions — generate **as many MCQs as required** to fully evaluate understanding of the content.
    - Each question must have **exactly 4 options** (labeled a, b, c, d).
    - Ensure **only one correct answer** per question.
    - Make sure the questions are diverse and non-repetitive.
    - The MCQs should reflect a mix of **factual, conceptual, and application-based understanding** of the content.

    **Content to generate MCQs from:**
    {text}

    **Output format (very important):**
    Q1. [Question]
    a) Option A  
    b) Option B  
    c) Option C  
    d) Option D  
    Answer: [Correct Option Letter]  
    """

    response = client.models.generate_content(
        model = "gemini-2.5-flash",
        contents = prompt
    )
    return response.text

# Streamlit UI
st.set_page_config(page_title="AI MCQ Generator", page_icon="📚", layout="wide")
st.title("📚 AI MCQ Generator using Gemini")

option = st.radio("Choose input method:", ("Text Input", "Upload Files", "YouTube Video Tutorial", "Recorded Lecture Audio"))

if option == "Text Input":
    user_text = st.text_area("Enter the text/content:")
    if st.button("Generate MCQs") and user_text.strip():
        with st.spinner("Generating MCQs..."):
            result = generate_mcqs(user_text)
            st.code(result,language="markdown")

elif option == "Upload Files":
    uploaded_files = st.file_uploader("Upload PDF, DOCX, or PPTX files", type=["pdf", "docx", "pptx"], accept_multiple_files=True)
    if st.button("Generate MCQs from Files") and uploaded_files:
        with st.spinner("Extracting and analyzing..."):
            extracted_text = extract_text_from_files(uploaded_files)
            if extracted_text.strip():
                result = generate_mcqs(extracted_text)
                st.code(result,language="markdown")
            else:
                st.warning("No text could be extracted from the uploaded files.")

elif option == "YouTube Video Tutorial":
    youtube_link = st.text_input("Enter YouTube Video URL:")

    def extract_transcript_details(youtube_video_url):
        try:
            video_id = youtube_video_url.split("=")[1]

            all_languages = [
                'af', 'sq', 'am', 'ar', 'hy', 'az', 'eu', 'be', 'bn', 'bs', 'bg', 'ca', 'ceb', 'zh', 'zh-Hans', 'zh-Hant',
                'co', 'hr', 'cs', 'da', 'nl', 'en', 'eo', 'et', 'fi', 'fr', 'fy', 'gl', 'ka', 'de', 'el', 'gu', 'ht', 'ha',
                'haw', 'he', 'hi', 'hmn', 'hu', 'is', 'ig', 'id', 'ga', 'it', 'ja', 'jv', 'kn', 'kk', 'km', 'rw', 'ko', 'ku',
                'ky', 'lo', 'la', 'lv', 'lt', 'lb', 'mk', 'mg', 'ms', 'ml', 'mt', 'mi', 'mr', 'mn', 'my', 'ne', 'no', 'ny',
                'or', 'ps', 'fa', 'pl', 'pt', 'pa', 'ro', 'ru', 'sm', 'gd', 'sr', 'st', 'sn', 'sd', 'si', 'sk', 'sl', 'so',
                'es', 'su', 'sw', 'sv', 'tl', 'tg', 'ta', 'tt', 'te', 'th', 'tr', 'tk', 'uk', 'ur', 'ug', 'uz', 'vi', 'cy',
                'xh', 'yi', 'yo', 'zu'
            ]

            transcript_text = YouTubeTranscriptApi.get_transcript(video_id, languages=all_languages)
            transcript = ""
            for i in transcript_text:
                transcript += " " + i["text"]
            return transcript
        except Exception as e:
            raise e

    if youtube_link:
        try:
            video_id = youtube_link.split("=")[1]
            st.image(f"http://img.youtube.com/vi/{video_id}/0.jpg", use_column_width=True)
        except:
            st.warning("Invalid YouTube URL format.")

    if st.button("Generate MCQs from YouTube Video"):
        with st.spinner("Extracting Transcript and Generating..."):
            transcript_text = extract_transcript_details(youtube_link)
            result = generate_mcqs(transcript_text)
            st.code(result,language="markdown")

elif option == "Recorded Lecture Audio":
    def transcribe_audio(audio_file_path):
        myfile = client.files.upload(file=audio_file_path)
        response = client.models.generate_content(
            model = "gemini-2.5-flash-preview-04-17",
            contents = ["Transcribe this complete audio clip",myfile]
        )
        return response.text.strip()

    def save_uploaded_file(uploaded_file):
        """Save uploaded file to a temporary file and return the path."""
        try:
            with tempfile.NamedTemporaryFile(delete=False, suffix='.' + uploaded_file.name.split('.')[-1]) as tmp_file:
                tmp_file.write(uploaded_file.getvalue())
                return tmp_file.name
        except Exception as e:
            st.error(f"Error handling uploaded file: {e}")
            return None

    audio_file = st.file_uploader("Upload Audio File", type=['wav', 'mp3'])
    if audio_file is not None:
        audio_path = save_uploaded_file(audio_file)  # Save the uploaded file and get the path
        st.audio(audio_path)

        if st.button('Transcribe and Generate MCQs'):
            with st.spinner('Transcribing and Generating...'):
                transcribed_text = transcribe_audio(audio_path)
                result = generate_mcqs(transcribed_text)
                st.code(result,language="markdown")
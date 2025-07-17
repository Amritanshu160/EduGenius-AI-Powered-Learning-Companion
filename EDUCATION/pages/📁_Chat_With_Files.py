import streamlit as st
import os
from PyPDF2 import PdfReader
from docx import Document
from pptx import Presentation
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain_google_genai import GoogleGenerativeAIEmbeddings
from google import genai
from langchain_community.vectorstores import FAISS
from langchain_google_genai import ChatGoogleGenerativeAI
from langchain.chains.question_answering import load_qa_chain
from langchain.prompts import PromptTemplate
from dotenv import load_dotenv

load_dotenv()

client = genai.Client(api_key=os.getenv("GOOGLE_API_KEY"))

def extract_text_from_files(files):
    text = ""
    for file in files:
        if file.name.endswith(".pdf"):
            pdf_reader = PdfReader(file)
            for page in pdf_reader.pages:
                text += page.extract_text()
        elif file.name.endswith(".docx"):
            doc = Document(file)
            for para in doc.paragraphs:
                text += para.text + "\n"
        elif file.name.endswith((".pptx", ".ppt")):
            presentation = Presentation(file)
            for slide in presentation.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
        else:
            st.warning(f"Unsupported file format: {file.name}")
    return text


def get_text_chunks(text):
    text_splitter = RecursiveCharacterTextSplitter(chunk_size=10000, chunk_overlap=1000)
    chunks = text_splitter.split_text(text)
    return chunks

def get_vector_store(text_chunks):
    embeddings = GoogleGenerativeAIEmbeddings(model="models/embedding-001")
    vector_store = FAISS.from_texts(text_chunks, embedding=embeddings)
    vector_store.save_local("faiss_index")

def get_conversational_chain():
    prompt_template = """
    Answer the question as detailed as possible from the provided context. Make sure to provide all the details. 
    If the answer is not in the provided context, just say, "Answer is not available in the context." 
    Don't provide a wrong answer.
    
    Context:\n{context}\n
    Question:\n{question}\n

    Answer:
    """
    model = ChatGoogleGenerativeAI(model="gemini-2.5-flash", temperature=0.3)
    prompt = PromptTemplate(template=prompt_template, input_variables=["context", "question"])
    chain = load_qa_chain(model, chain_type="stuff", prompt=prompt)
    return chain

def generate_response(user_question):
    embeddings = GoogleGenerativeAIEmbeddings(model="models/embedding-001")
    new_db = FAISS.load_local("faiss_index", embeddings, allow_dangerous_deserialization=True)
    docs = new_db.similarity_search(user_question)
    chain = get_conversational_chain()

    response = chain(
        {"input_documents": docs, "question": user_question},
        return_only_outputs=True
    )

    return response["output_text"]

def main():
    st.set_page_config(page_title="Chat With Multiple Files",page_icon="📁",layout="wide")
    st.title("Chat with Multiple Files using Gemini 📄✨")

    if "messages" not in st.session_state:
        st.session_state["messages"] = [
            {"role": "assistant", "content": "Hello! Upload your files from the sidebar and ask me anything from them!"}
        ]

    # Display the chat history
    for msg in st.session_state.messages:
        st.chat_message(msg["role"]).write(msg["content"])

    # Chat input
    if prompt := st.chat_input("Ask a question about your uploads..."):
        st.session_state.messages.append({"role": "user", "content": prompt})
        st.chat_message("user").write(prompt)

        with st.chat_message("assistant"):
            with st.spinner("Thinking..."):
                try:
                    response = generate_response(prompt)
                except Exception as e:
                    response = f"Error: {e}"

            st.session_state.messages.append({"role": "assistant", "content": response})
            st.write(response)

    # Sidebar for uploading PDFs
    with st.sidebar:
        st.title("Menu:")
        uploaded_files = st.file_uploader("Upload your Files (PDF, Word, PPT) and Click on Submit & Process", accept_multiple_files=True, type=['pdf', 'docx', 'pptx', 'ppt'])
        if st.button("Submit & Process"):
            if uploaded_files:
                with st.spinner("Processing..."):
                    raw_text = extract_text_from_files(uploaded_files)
                    text_chunks = get_text_chunks(raw_text)
                    get_vector_store(text_chunks)
                    st.success("Files processed successfully! 🎉")
            else:
                st.warning("Please upload PDF files before submitting.")

if __name__ == "__main__":
    main()
